import base64

import pptx
from pptx.util import Inches, Pt
import os
from datetime import datetime

 
# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

 

# Function to get response
def get_response(user_message):
    

    return "response.choices[0].message.content"
def generate_slide_titles(topic):


    prompt = f"Generate 3 slide titles for the topic '{topic}'."

    response = get_response(prompt)
    print(response)

    return response
 

def generate_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'."
   
    return prompt
 
def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

    if not os.path.exists('ppts'):
    os.makedirs('ppts')
    now = datetime.now().strftime("%d_%m_%Y_%H_%M")
    prs.save(f"ppts/{topic}_presentation_{now}.pptx")
    #prs.save(f"{topic}_presentation.pptx")


def get_ppt_download_link(topic):
    ppt_filename = f"{topic}_presentation.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}'


def generateppt(topic):

    

    slide_titles = generate_slide_titles(topic)
    filtered_slide_titles= [item for item in slide_titles if item.strip() != '']
    print("Slide Title: ", filtered_slide_titles)
    slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
    print("Slide Contents: ", slide_contents)
    create_presentation(topic, filtered_slide_titles, slide_contents)
     

    #print("Total Duration: ", round(end - start, 2), " secs")
    print("Presentation generated successfully!")

    download_link = get_ppt_download_link(topic)

    return download_link





